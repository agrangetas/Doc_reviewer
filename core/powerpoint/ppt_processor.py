"""
PowerPoint Processor - Implémentation complète
Traitement de présentations PowerPoint avec IA.
"""

from pathlib import Path
from typing import Optional
from pptx import Presentation
from datetime import datetime

from core.base.document_processor import DocumentProcessor


class PowerPointProcessor(DocumentProcessor):
    """Processeur pour les présentations PowerPoint (.pptx)."""
    
    def __init__(self, config, image_handler, style_extractor, style_mapper,
                 language_detector, ai_processor, logger, style_uniformizer):
        """Initialise le processeur PowerPoint."""
        super().__init__(
            config, image_handler, style_extractor, style_mapper,
            language_detector, ai_processor, logger, style_uniformizer
        )
        self.presentation = None
        self.cached_structure = None  # Cache pour la structure des slides
    
    def load_document(self, file_path: str) -> None:
        """
        Charge une présentation PowerPoint.
        
        Args:
            file_path: Chemin vers le fichier .pptx
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Le fichier {file_path} n'existe pas.")
        
        self.presentation = Presentation(file_path)
        self.current_path = path
        
        # Détecter la langue
        sample_text = self._extract_sample_text()
        if sample_text:
            lang_code = self.language_detector.detect_language(sample_text)
            if lang_code:
                self.detected_language = lang_code
        
        # Calculer la structure UNE SEULE FOIS au chargement
        print("📊 Analyse de la structure de la présentation...")
        self.cached_structure = self._calculate_structure()
        
        # Afficher les infos
        print(f"✓ Présentation chargée: {path.name}")
        print(f"  Nombre de slides: {len(self.presentation.slides)}")
        print(f"  Modèle OpenAI: {self.ai_processor.model}")
        if self.detected_language:
            print(f"  Langue détectée: {self.language_detector.get_language_name(self.detected_language)}")
        
        # Initialiser le log
        doc_info = {
            'slide_count': len(self.presentation.slides),
            'language': self.language_detector.get_language_name(self.detected_language) if self.detected_language else None
        }
        self.logger.init_log_file(path.name, doc_info)
    
    def _extract_sample_text(self) -> str:
        """Extrait un échantillon de texte pour détection de langue."""
        sample = []
        for i, slide in enumerate(self.presentation.slides):
            if i >= 3:  # Limiter à 3 slides
                break
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    sample.append(shape.text.strip())
        return " ".join(sample)
    
    def _preserve_shape_format(self, shape, new_text: str) -> None:
        """
        Préserve le formatage d'un shape PowerPoint en appliquant le nouveau texte.
        Utilise le système de mapping de styles comme pour Word.
        
        Args:
            shape: Shape PowerPoint
            new_text: Nouveau texte à appliquer
        """
        if not hasattr(shape, 'text_frame'):
            shape.text = new_text
            return
        
        text_frame = shape.text_frame
        
        # Sauvegarder les propriétés des paragraphes (alignement, bullet points, etc.)
        paragraph_formats = []
        for para in text_frame.paragraphs:
            para_format = {
                'alignment': para.alignment,
                'level': para.level,
                'line_spacing': para.line_spacing,
                'space_before': para.space_before,
                'space_after': para.space_after,
            }
            # Bullet points
            if hasattr(para, 'font') and para.font:
                para_format['bullet_font'] = {
                    'name': para.font.name,
                    'size': para.font.size,
                    'bold': para.font.bold,
                }
            paragraph_formats.append(para_format)
        
        # Extraire le texte original complet
        original_text = shape.text
        
        # Extraire les styles avec le système existant (adapté pour PPT)
        styles_map = self._extract_ppt_styles_map(text_frame)
        
        # Mapper les styles vers le nouveau texte
        new_styles_map = self.style_mapper.map_styles_to_new_text(
            original_text, new_text, styles_map
        )
        
        # Appliquer le nouveau texte avec les styles mappés
        self._apply_ppt_styles_map(text_frame, new_text, new_styles_map, paragraph_formats)
    
    def _extract_ppt_styles_map(self, text_frame):
        """
        Extrait la map des styles d'un text_frame PowerPoint.
        Similaire à StyleExtractor mais adapté pour PPT.
        
        Args:
            text_frame: TextFrame PowerPoint
            
        Returns:
            Liste de dicts avec positions et styles
        """
        styles_map = []
        current_position = 0
        
        for para in text_frame.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue
                
                run_length = len(run.text)
                
                styles_map.append({
                    'start': current_position,
                    'end': current_position + run_length,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'underline': run.font.underline,
                    'font_name': run.font.name,
                    'font_size': run.font.size,
                    'font_color': run.font.color.rgb if hasattr(run.font.color, 'rgb') and run.font.color.type == 1 else None
                })
                
                current_position += run_length
        
        return styles_map
    
    def _apply_ppt_styles_map(self, text_frame, new_text: str, styles_map: list, paragraph_formats: list):
        """
        Applique une map de styles à un text_frame PowerPoint.
        
        Args:
            text_frame: TextFrame PowerPoint
            new_text: Nouveau texte
            styles_map: Map des styles
            paragraph_formats: Formats des paragraphes (alignement, bullets, etc.)
        """
        # Nettoyer le text_frame
        text_frame.clear()
        
        # Si pas de styles, juste ajouter le texte avec le format du premier paragraphe
        if not styles_map:
            p = text_frame.paragraphs[0]
            p.text = new_text
            if paragraph_formats:
                self._apply_paragraph_format(p, paragraph_formats[0])
            return
        
        # Créer un paragraphe et ajouter les runs avec styles
        p = text_frame.paragraphs[0]
        
        # Appliquer le format du premier paragraphe (alignement, bullets)
        if paragraph_formats:
            self._apply_paragraph_format(p, paragraph_formats[0])
        
        # Trier les styles par position
        sorted_styles = sorted(styles_map, key=lambda x: x['start'])
        
        # Créer des runs pour chaque section de style
        last_end = 0
        
        for style in sorted_styles:
            start = max(style['start'], last_end)
            end = min(style['end'], len(new_text))
            
            if start >= len(new_text):
                break
            
            # Texte avant ce style (si gap)
            if start > last_end:
                gap_text = new_text[last_end:start]
                if gap_text:
                    p.add_run().text = gap_text
            
            # Texte avec ce style
            if end > start:
                styled_text = new_text[start:end]
                run = p.add_run()
                run.text = styled_text
                
                # Appliquer le style
                if style['bold'] is not None:
                    run.font.bold = style['bold']
                if style['italic'] is not None:
                    run.font.italic = style['italic']
                if style['underline'] is not None:
                    run.font.underline = style['underline']
                if style['font_name']:
                    run.font.name = style['font_name']
                if style['font_size']:
                    run.font.size = style['font_size']
                if style['font_color']:
                    try:
                        run.font.color.rgb = style['font_color']
                    except:
                        pass
                
                last_end = end
        
        # Texte restant après tous les styles
        if last_end < len(new_text):
            remaining_text = new_text[last_end:]
            if remaining_text:
                p.add_run().text = remaining_text
    
    def _apply_paragraph_format(self, paragraph, para_format: dict):
        """
        Applique le formatage de paragraphe (alignement, bullets, etc.).
        
        Args:
            paragraph: Paragraphe PowerPoint
            para_format: Dict avec les formats à appliquer
        """
        try:
            if para_format.get('alignment') is not None:
                paragraph.alignment = para_format['alignment']
            if para_format.get('level') is not None:
                paragraph.level = para_format['level']
            if para_format.get('line_spacing') is not None:
                paragraph.line_spacing = para_format['line_spacing']
            if para_format.get('space_before') is not None:
                paragraph.space_before = para_format['space_before']
            if para_format.get('space_after') is not None:
                paragraph.space_after = para_format['space_after']
        except Exception as e:
            # Certains formats peuvent ne pas être applicables
            pass
    
    def _calculate_structure(self):
        """
        Calcule la structure de la présentation (une seule fois).
        
        Returns:
            Dict avec informations sur les slides et shapes
        """
        structure = {
            'total_slides': len(self.presentation.slides),
            'slides': []
        }
        
        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            slide_info = {
                'number': slide_idx,
                'shape_count': len([s for s in slide.shapes if hasattr(s, "text")]),
                'has_title': False,
                'has_images': False
            }
            
            # Détecter titre et images
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Vérifier si c'est probablement un titre (position haute, texte court)
                    try:
                        if shape.top is not None and shape.top < 2000000 and len(shape.text) < 100:
                            slide_info['has_title'] = True
                    except:
                        pass
                
                try:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        slide_info['has_images'] = True
                except:
                    pass
            
            structure['slides'].append(slide_info)
        
        print(f"   ✓ Structure analysée: {structure['total_slides']} slides")
        return structure
    
    def save_document(self, output_path: Optional[str] = None) -> None:
        """
        Sauvegarde la présentation PowerPoint.
        
        Args:
            output_path: Chemin de sortie (optionnel)
        """
        if not self.presentation:
            raise ValueError("Aucune présentation à sauvegarder.")
        
        if output_path is None:
            output_path = self.current_path.parent / f"{self.current_path.stem}_modifié{self.current_path.suffix}"
        
        self.presentation.save(output_path)
        print(f"\n💾 Présentation sauvegardée: {output_path}")
    
    def process_document(self, instruction: str) -> None:
        """
        Traite la présentation avec une instruction.
        
        Args:
            instruction: Instruction à exécuter
        """
        if not self.presentation:
            raise ValueError("Aucune présentation chargée.")
        
        is_correction = any(word in instruction.lower() for word in ['corrige', 'correction', 'orthographe', 'grammaire'])
        
        from features.language_detector import LanguageDetector
        language_name = LanguageDetector.get_language_name(self.detected_language) if self.detected_language else None
        
        print(f"\n{'=' * 60}")
        print(f"TRAITEMENT : {instruction}")
        if is_correction and language_name:
            print(f"Langue détectée : {language_name}")
        print('=' * 60)
        
        total_shapes = sum(len([s for s in slide.shapes if hasattr(s, "text")]) for slide in self.presentation.slides)
        current_shape = 0
        
        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            print(f"\nSlide {slide_idx}/{len(self.presentation.slides)}...")
            
            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                
                current_shape += 1
                original_text = shape.text
                
                print(f"  Shape {current_shape}/{total_shapes}...", end=" ", flush=True)
                
                try:
                    # Construire le contexte
                    context_parts = []
                    if is_correction and language_name:
                        context_parts.append(f"Le texte est en {language_name}.")
                    
                    context = " ".join(context_parts) if context_parts else ""
                    
                    # Appeler l'IA
                    new_text = self.ai_processor.call_openai(
                        instruction, original_text, context, is_correction, language_name
                    )
                    
                    if new_text and new_text.strip() != original_text.strip():
                        # Nettoyer le texte retourné (enlever "text:" si présent)
                        if new_text.lower().startswith("text:"):
                            new_text = new_text[5:].strip()
                        
                        # Préserver le formatage
                        self._preserve_shape_format(shape, new_text)
                        print("✓ Modifié")
                        
                        # Logger
                        self.logger.log_change(
                            f'slide_{slide_idx}_shape_{current_shape}',
                            original_text,
                            new_text,
                            instruction
                        )
                    else:
                        print("○ Inchangé")
                
                except Exception as e:
                    print(f"❌ Erreur: {e}")
        
        print(f"\n{'=' * 60}")
        print("✓ Traitement terminé")
        print('=' * 60)
    
    def process_targeted(self, target: dict, instruction: str) -> None:
        """
        Traite un élément spécifique de la présentation.
        
        Args:
            target: Informations sur l'élément ciblé (depuis ElementResolver)
            instruction: Instruction à exécuter
        """
        if not self.presentation:
            raise ValueError("Aucune présentation chargée.")
        
        slide_number = target.slide
        shape_index = target.shape
        
        if not slide_number or slide_number < 1 or slide_number > len(self.presentation.slides):
            raise ValueError(f"Numéro de slide invalide: {slide_number}")
        
        slide = self.presentation.slides[slide_number - 1]
        
        print(f"\n{'=' * 60}")
        print(f"TRAITEMENT CIBLÉ 🎯")
        print(f"Slide: {slide_number}")
        print(f"Instruction: {instruction}")
        print('=' * 60)
        
        # Trouver le shape ciblé
        if shape_index is not None:
            shapes_with_text = [s for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if shape_index < len(shapes_with_text):
                shape = shapes_with_text[shape_index]
                original_text = shape.text
                
                print(f"\n📝 Texte original:\n{original_text}\n")
                
                try:
                    # Construire le contexte
                    context_parts = []
                    is_correction = any(word in instruction.lower() for word in ['corrige', 'correction', 'orthographe', 'grammaire'])
                    lang_name = None
                    if self.detected_language:
                        lang_name = self.language_detector.get_language_name(self.detected_language)
                        if is_correction:
                            context_parts.append(f"Le texte est en {lang_name}.")
                    
                    context = " ".join(context_parts) if context_parts else ""
                    
                    # Appeler l'IA
                    new_text = self.ai_processor.call_openai(
                        instruction, original_text, context, is_correction, lang_name
                    )
                    
                    if new_text and new_text.strip() != original_text.strip():
                        # Nettoyer le texte retourné (enlever "text:" si présent)
                        if new_text.lower().startswith("text:"):
                            new_text = new_text[5:].strip()
                        
                        # Préserver le formatage
                        self._preserve_shape_format(shape, new_text)
                        print(f"✅ Texte modifié:\n{new_text}\n")
                        
                        # Logger
                        self.logger.log_change(
                            f'slide_{slide_number}_shape_{shape_index}',
                            original_text,
                            new_text,
                            f"{instruction} (ciblé)"
                        )
                    else:
                        print("○ Aucune modification nécessaire")
                
                except Exception as e:
                    print(f"❌ Erreur: {e}")
            else:
                print(f"⚠️  Shape index invalide: {shape_index}")
        
        print(f"{'=' * 60}")
    
    def list_elements(self) -> None:
        """Liste tous les éléments de la présentation."""
        if not self.presentation:
            raise ValueError("Aucune présentation chargée.")
        
        print(f"\n{'=' * 60}")
        print("STRUCTURE DE LA PRÉSENTATION")
        print('=' * 60)
        
        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            print(f"\nSlide {slide_idx}:")
            shape_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shape_count += 1
                    preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                    print(f"  [{shape_count}] {preview}")
        
        print(f"\n{'=' * 60}")
    
    def uniformize_styles(self) -> None:
        """Uniformise les styles de la présentation."""
        print("\n⚠️  L'uniformisation de styles PowerPoint n'est pas encore implémentée.")
        print("Cette fonctionnalité sera ajoutée dans une future version.")
    
    def get_format_name(self) -> str:
        """Retourne le nom du format du document."""
        return "PowerPoint"

